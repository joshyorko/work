from pptx import Presentation

# Path to the template PowerPoint file
template_path = "path/to/template.pptx"

# Information for the slide
slide_title = "Retention Policy Summary"
content_info = [
    ("Sharepoint", "No retention"),
    ("Teams chat", "60 days"),
    ("Teams channels", "No retention"),
    ("OneDrive", "To be confirmed")
]

# Load the template presentation
presentation = Presentation(template_path)

# Create a new slide
slide_layout = presentation.slide_layouts[1]  # Choose the layout for a content slide
slide = presentation.slides.add_slide(slide_layout)

# Set the slide title
title_placeholder = slide.shapes.title
title_placeholder.text = slide_title

# Add content information to the slide
content_placeholder = slide.placeholders[1]
for item in content_info:
    content = f"{item[0]}: {item[1]}"
    p = content_placeholder.text_frame.add_paragraph()
    p.text = content

# Save the generated presentation
output_path = "path/to/output.pptx"
presentation.save(output_path)
